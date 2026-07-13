from __future__ import annotations

from pathlib import Path

from services.parsers.base import ParsedDocument, format_table


def parse_pptx(path: str | Path) -> ParsedDocument:
    """Парсер презентаций (.pptx). Каждый слайд → блок «Слайд N», внутри — текст
    фигур (заголовки, списки), таблицы (с сохранением связи строка↔столбец) и
    заметки докладчика. Подходит для схем процессов, оформленных в PowerPoint."""
    path = Path(path)
    try:
        from pptx import Presentation
    except ImportError as e:  # pragma: no cover
        raise ValueError(
            "Для .pptx нужен пакет python-pptx (pip install python-pptx)."
        ) from e

    prs = Presentation(str(path))
    parts: list[str] = []

    for idx, slide in enumerate(prs.slides, start=1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            # Таблицы размечаем как в docx/xlsx — значение не теряет смысл при чанковании.
            if shape.has_table:
                rows = [
                    [(cell.text or "").strip() for cell in row.cells]
                    for row in shape.table.rows
                ]
                slide_parts.extend(format_table(rows))
                continue
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        slide_parts.append(text)

        # Заметки докладчика — часто содержат пояснения к схеме.
        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
            if notes:
                slide_parts.append(f"Заметки: {notes}")

        if slide_parts:
            parts.append(f"Слайд {idx}")
            parts.extend(slide_parts)

    return ParsedDocument(
        text="\n\n".join(parts),
        title=path.stem,
        source_uri=str(path),
        source_type="local",
        mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
